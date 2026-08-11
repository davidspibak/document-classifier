"""
Format-specific parsers. Each parser returns a dict:
    {"text": str, "pages": list[str] | None, "needs_ocr_pages": list[int]}
`pages` is only populated for PDFs (needed by ocr.py to know which page images
to extract); other formats just return the full text.
"""
from pathlib import Path
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from bs4 import BeautifulSoup

# A page with fewer than this many characters of extracted text is treated as
# "no text layer" and queued for OCR instead of trusted as-is.
MIN_TEXT_LENGTH_PER_PAGE = 20


def parse_pdf(path: str) -> dict:
    pages_text = []
    needs_ocr_pages = []
    # Context-managed so the file handle is released before ocr.py re-opens the
    # same path — a bulk run that leaks one handle per document hits the OS limit.
    with fitz.open(path) as doc:
        for i, page in enumerate(doc):
            text = page.get_text().strip()
            pages_text.append(text)
            if len(text) < MIN_TEXT_LENGTH_PER_PAGE:
                needs_ocr_pages.append(i)
    return {
        "text": "\n\n".join(pages_text),
        "pages": pages_text,
        "needs_ocr_pages": needs_ocr_pages,
        "pdf_handle_path": path,  # ocr.py re-opens the file to rasterize flagged pages
    }


def parse_docx(path: str) -> dict:
    d = docx.Document(path)
    paragraphs = [p.text for p in d.paragraphs if p.text.strip()]
    # Tables often carry real content (e.g. data tables in reports) — include them.
    for table in d.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip(" |"):
                paragraphs.append(row_text)
    return {"text": "\n".join(paragraphs), "pages": None, "needs_ocr_pages": []}


def parse_pptx(path: str) -> dict:
    prs = Presentation(path)
    slides_text = []
    for slide in prs.slides:
        parts = [shape.text for shape in slide.shapes if shape.has_text_frame]
        slides_text.append("\n".join(parts))
    return {"text": "\n\n".join(slides_text), "pages": slides_text, "needs_ocr_pages": []}


def parse_html(path: str) -> dict:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    # Strip boilerplate that would otherwise pollute embeddings/classification.
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    return {"text": soup.get_text(separator="\n").strip(), "pages": None, "needs_ocr_pages": []}


def parse_txt(path: str) -> dict:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return {"text": f.read(), "pages": None, "needs_ocr_pages": []}


_PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
    ".html": parse_html,
    ".htm": parse_html,
    ".txt": parse_txt,
    ".md": parse_txt,
}


def parse_document(path: str) -> dict:
    """Dispatches to the right parser based on file extension."""
    ext = Path(path).suffix.lower()
    parser = _PARSERS.get(ext)
    if parser is None:
        raise ValueError(f"No parser registered for extension '{ext}' (file: {path})")
    return parser(path)
